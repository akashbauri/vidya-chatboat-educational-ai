import streamlit as st
import os
import google.generativeai as genai
import faiss
import numpy as np
from sentence_transformers import SentenceTransformer
import PyPDF2
from pptx import Presentation
import requests
from bs4 import BeautifulSoup
from duckduckgo_search import DDGS
import json
from datetime import datetime
import re

# ============================================================
# PAGE CONFIGURATION
# ============================================================
# Configure the Streamlit page with title, icon, and layout
st.set_page_config(
    page_title="Vidya Chatbot - AI Study Assistant",
    page_icon="🧩",
    layout="wide"
)

# ============================================================
# GEMINI API SETUP
# ============================================================
# Initialize Google's Gemini AI model for generating responses
try:
    # Configure API key from Streamlit secrets
    genai.configure(api_key=st.secrets["GEMINI_API_KEY"])
    
    # Get SDK version and available models for debugging
    sdk_version = genai.__version__
    available_models = [m.name for m in genai.list_models()]

    # Display model info in sidebar
    st.sidebar.write("📦 Gemini SDK Version:", sdk_version)
    st.sidebar.write("✅ Available Models:", available_models)

    # List of preferred Gemini models in order of preference
    model_name = None
    preferred_order = [
        "models/gemini-2.5-flash",
        "models/gemini-2.5-flash-preview-05-20",
        "models/gemini-2.5-flash-lite-preview-06-17",
        "models/gemini-2.5-pro-preview-05-06",
        "models/gemini-2.5-pro-preview-03-25",
    ]

    # Try to find and use the first available preferred model
    for name in preferred_order:
        if name in available_models:
            model_name = name
            break

    # If preferred model found, initialize it
    if model_name:
        model = genai.GenerativeModel(model_name)
        st.sidebar.success(f"✅ Using Model: {model_name}")
    else:
        # If no preferred model, try to use any available Gemini model
        for m in available_models:
            if "gemini" in m:
                model_name = m
                model = genai.GenerativeModel(model_name)
                st.sidebar.warning(f"⚙️ Using fallback model: {model_name}")
                break
        if not model_name:
            st.error("⚠️ No supported Gemini model found. Please check your API key.")
            model = None

except Exception as e:
    st.error(f"⚠️ Gemini initialization failed: {e}")
    model = None

# ============================================================
# EMBEDDING MODEL SETUP
# ============================================================
# Load sentence-transformers model for converting text to embeddings
# This model helps us find similar text chunks for RAG
@st.cache_resource
def load_embedding_model():
    """
    Load and cache the sentence transformer model.
    This model converts text into numerical vectors (embeddings).
    We use 'all-MiniLM-L6-v2' because it's fast and accurate.
    """
    return SentenceTransformer('sentence-transformers/all-MiniLM-L6-v2')

embedding_model = load_embedding_model()

# ============================================================
# SESSION STATE INITIALIZATION
# ============================================================
# Initialize session state variables to store data across reruns
for key in ['messages', 'vector_store', 'documents', 'embeddings', 'sources']:
    if key not in st.session_state:
        # Messages should be a list, others start as None
        st.session_state[key] = [] if key == 'messages' else None

# ============================================================
# CONFIGURATION CONSTANTS
# ============================================================
# Define app-wide constants for file processing and retrieval
MAX_FILE_SIZE = 500 * 1024 * 1024  # Maximum file size: 500 MB
CHUNK_SIZE = 1000                  # Number of characters per text chunk
CHUNK_OVERLAP = 150                # Overlap between consecutive chunks
SIMILARITY_THRESHOLD = 0.25        # Minimum similarity score to consider a chunk relevant
TOP_K = 6                          # Number of top similar chunks to retrieve

# ============================================================
# TEXT EXTRACTION FUNCTIONS
# ============================================================

def extract_text_from_pdf(file):
    """
    Extract text from PDF file and add source information.
    Each page's text is tagged with filename and page number.
    
    Args:
        file: Uploaded PDF file object
        
    Returns:
        String containing all text with source tags, or None if error
    """
    try:
        pdf_reader = PyPDF2.PdfReader(file)
        text = ""
        
        # Loop through each page and extract text
        for i, page in enumerate(pdf_reader.pages):
            page_text = page.extract_text()
            if page_text:
                # Add source tag before each page's content
                text += f"\n[Source: {file.name}, Page {i+1}]\n{page_text}"
        
        # Return None if no text was extracted
        return text or None
        
    except Exception as e:
        st.error(f"Error reading PDF: {e}")
        return None

def extract_text_from_pptx(file):
    """
    Extract text from PowerPoint file and add source information.
    Each slide's text is tagged with filename and slide number.
    
    Args:
        file: Uploaded PPTX file object
        
    Returns:
        String containing all text with source tags, or None if error
    """
    try:
        prs = Presentation(file)
        text = ""
        
        # Loop through each slide
        for i, slide in enumerate(prs.slides):
            slide_text = ""
            
            # Extract text from all shapes in the slide
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    slide_text += shape.text + "\n"
            
            # Add source tag if slide has content
            if slide_text.strip():
                text += f"\n[Source: {file.name}, Slide {i+1}]\n{slide_text}"
        
        return text or None
        
    except Exception as e:
        st.error(f"Error reading PPTX: {e}")
        return None

def extract_text_from_url(url):
    """
    Fetch and extract clean text from a webpage.
    Removes scripts, styles, and extra whitespace.
    
    Args:
        url: Web address to fetch
        
    Returns:
        String containing cleaned text with source tag, or None if error
    """
    try:
        # Set user agent to avoid being blocked
        headers = {'User-Agent': 'Mozilla/5.0'}
        response = requests.get(url, headers=headers, timeout=10)
        response.raise_for_status()
        
        # Parse HTML content
        soup = BeautifulSoup(response.content, 'html.parser')
        
        # Remove script and style elements
        for s in soup(["script", "style"]):
            s.decompose()
        
        # Extract and clean visible text
        text = " ".join(t.strip() for t in soup.get_text().splitlines() if t.strip())
        
        return f"[Source: {url}]\n{text}" if text else None
        
    except Exception as e:
        st.error(f"Error fetching URL: {e}")
        return None

# ============================================================
# TEXT CHUNKING FUNCTION
# ============================================================

def chunk_text(text, chunk_size=CHUNK_SIZE, overlap=CHUNK_OVERLAP):
    """
    Split large text into smaller overlapping chunks.
    Overlapping helps maintain context across chunk boundaries.
    
    Args:
        text: String to split into chunks
        chunk_size: Maximum characters per chunk
        overlap: Number of characters to overlap between chunks
        
    Returns:
        List of text chunks
    """
    if not text:
        return []
    
    chunks = []
    start = 0
    text_length = len(text)
    
    # Create chunks with sliding window approach
    while start < text_length:
        end = start + chunk_size
        chunk = text[start:end]
        
        # Only add non-empty chunks
        if chunk.strip():
            chunks.append(chunk)
        
        # Move window forward (chunk_size - overlap)
        # Minimum step is 1 to avoid infinite loop
        start += max(1, chunk_size - overlap)
    
    return chunks

# ============================================================
# VECTOR STORE FUNCTIONS (FAISS + EMBEDDINGS)
# ============================================================

def create_vector_store(documents):
    """
    Create FAISS vector index from document chunks.
    This function does the following:
    1. Cleans all documents to ensure they are plain strings
    2. Converts strings to embeddings using sentence-transformers
    3. Normalizes embeddings for cosine similarity
    4. Creates FAISS index for fast similarity search
    
    Args:
        documents: List of text chunks
        
    Returns:
        Tuple of (FAISS index, embeddings array) or (None, None) if no valid docs
    """
    # Return early if no documents provided
    if not documents:
        return None, None

    # CRITICAL CLEANING STEP to fix the TextEncodeInput error
    # sentence-transformers requires a list of plain strings
    cleaned_docs = []
    
    for d in documents:
        # Skip None values
        if d is None:
            continue
        
        # If somehow a list or tuple got in, join it into a string
        if isinstance(d, (list, tuple)):
            d = " ".join(str(x) for x in d)
        else:
            # Convert to string (handles any other type)
            d = str(d)
        
        # Remove leading/trailing whitespace
        d = d.strip()
        
        # Only keep non-empty strings
        if d:
            cleaned_docs.append(d)

    # If all documents were empty, return None
    if not cleaned_docs:
        return None, None

    # Create embeddings and FAISS index
    with st.spinner("Creating embeddings..."):
        # Convert text to embeddings (numerical vectors)
        embeddings = np.array(
            embedding_model.encode(cleaned_docs, show_progress_bar=True)
        ).astype('float32')
        
        # Normalize embeddings for cosine similarity
        faiss.normalize_L2(embeddings)
        
        # Create FAISS index for inner product (equivalent to cosine similarity after normalization)
        index = faiss.IndexFlatIP(embeddings.shape[1])
        
        # Add all embeddings to the index
        index.add(embeddings)
    
    return index, embeddings

def clean_text_for_embedding(text):
    """
    Clean text before converting to embedding.
    Removes punctuation and normalizes whitespace.
    
    Args:
        text: String to clean
        
    Returns:
        Cleaned string
    """
    # Remove all non-alphanumeric characters except spaces
    text = re.sub(r'[^\w\s]', '', text)
    
    # Replace multiple spaces with single space
    text = re.sub(r'\s+', ' ', text)
    
    return text.strip()

def retrieve_relevant_chunks(query, top_k=TOP_K):
    """
    Find the most relevant document chunks for a query.
    Uses FAISS to search for similar embeddings.
    
    Args:
        query: User's question
        top_k: Number of top results to return
        
    Returns:
        List of dictionaries with 'text' and 'similarity' keys
    """
    # Check if vector store exists
    if st.session_state.vector_store is None:
        return []
    
    # Clean the query text
    query_clean = clean_text_for_embedding(query)
    if not query_clean:
        return []
    
    # Convert query to embedding
    query_embedding = np.array(
        embedding_model.encode([query_clean])
    ).astype('float32')
    
    # Normalize query embedding
    faiss.normalize_L2(query_embedding)
    
    # Search FAISS index for similar chunks
    distances, indices = st.session_state.vector_store.search(query_embedding, top_k)
    
    # Filter results by similarity threshold
    results = []
    for d, i in zip(distances[0], indices[0]):
        # Only include chunks above similarity threshold
        # Also check index is valid
        if d >= SIMILARITY_THRESHOLD and 0 <= i < len(st.session_state.documents):
            results.append({
                'text': st.session_state.documents[i],
                'similarity': float(d)
            })
    
    return results

# ============================================================
# WEB SEARCH FUNCTION
# ============================================================

def web_search(query, num_results=3):
    """
    Search the web using DuckDuckGo when no relevant documents found.
    Provides fallback information when user materials don't have answers.
    
    Args:
        query: Search query
        num_results: Number of results to return
        
    Returns:
        List of search result dictionaries
    """
    try:
        with DDGS() as ddgs:
            # Perform web search
            return list(ddgs.text(query, max_results=num_results))
    except Exception as e:
        st.error(f"Web search error: {e}")
        return []

# ============================================================
# GEMINI RESPONSE GENERATION
# ============================================================

def generate_response(query, context_chunks=None, web_results=None):
    """
    Generate response using Gemini AI based on context.
    Three modes:
    1. RAG mode: Use retrieved document chunks
    2. Web mode: Use web search results
    3. Direct mode: Answer without context
    
    Args:
        query: User's question
        context_chunks: Retrieved document chunks (RAG mode)
        web_results: Web search results (Web mode)
        
    Returns:
        Generated response string
    """
    # Check if model is initialized
    if not model:
        return "⚠️ Gemini model not initialized properly."

    # RAG MODE: Use user's uploaded materials
    if context_chunks:
        # Build context from retrieved chunks
        context = "\n\n".join(
            f"[Context {i+1}] {chunk['text'][:500]}" 
            for i, chunk in enumerate(context_chunks)
        )
        
        # Create prompt with context
        prompt = f"""
You are Vidya, an intelligent educational assistant.
Answer the question using the context below and cite the file/slide/page source at the end of each point.

Context from user materials:
{context}

Question: {query}

Instructions:
- Always include sources in brackets (e.g., [Source: Python_Notes.pdf, Page 4])
- If information comes from multiple files, summarize them clearly
- Keep it short and clear for students
"""
    
    # WEB MODE: Use web search results
    elif web_results:
        # Build context from web search
        web_context = "\n\n".join(
            f"[Web {i+1}] {r['title']} - {r['body'][:300]} (Source: {r['href']})"
            for i, r in enumerate(web_results)
        )
        
        # Create prompt with web context
        prompt = f"""
You are Vidya, a helpful educational assistant.
Use the web search results below to answer the question accurately and include citations.

Web Results:
{web_context}

Question: {query}

Instructions:
- Summarize clearly using the most relevant sources
- Always include [Source: website.com] citations
- Keep your tone simple and educational
"""
    
    # DIRECT MODE: No context available
    else:
        prompt = f"You are Vidya, an educational AI. Answer this: {query}"

    # Generate response using Gemini
    try:
        response = model.generate_content(prompt)
        return response.text
    except Exception as e:
        st.error(f"⚠️ Error generating response: {e}")
        return "I encountered an issue generating a response."

# ============================================================
# STREAMLIT UI - MAIN SECTION
# ============================================================

# Display main title and description
st.title("🧩 Vidya Chatbot")
st.markdown("### AI-Powered Educational Assistant")
st.markdown("*Ask questions about your study materials or any topic!*")

# ============================================================
# SIDEBAR - FILE UPLOAD AND PROCESSING
# ============================================================

with st.sidebar:
    st.header("📚 Upload Study Materials")
    
    # File uploader widget
    uploaded_files = st.file_uploader(
        "Upload PDF or PPTX files", 
        type=['pdf', 'pptx'], 
        accept_multiple_files=True
    )
    
    # URL input widget
    url_input = st.text_input("Or enter a URL:")

    # Process materials button
    if st.button("Process Materials", type="primary"):
        all_text = []
        
        # Process each uploaded file
        for file in uploaded_files or []:
            # Check file size limit
            if file.size > MAX_FILE_SIZE:
                st.error(f"❌ {file.name} exceeds 500 MB limit")
                continue
            
            with st.spinner(f"Processing {file.name}..."):
                # Extract text based on file type
                if file.name.lower().endswith(".pdf"):
                    text = extract_text_from_pdf(file)
                else:
                    text = extract_text_from_pptx(file)
                
                # Add to collection if extraction successful
                if text:
                    all_text.append(text)
                    st.success(f"✅ {file.name} processed")

        # Process URL if provided
        if url_input:
            with st.spinner(f"Fetching {url_input}..."):
                text = extract_text_from_url(url_input)
                if text:
                    all_text.append(text)
                    st.success("✅ URL processed")

        # If we have any extracted text, create vector store
        if all_text:
            # Combine all extracted texts
            combined_text = "\n\n".join(all_text)
            
            # Split into chunks
            chunks = chunk_text(combined_text)
            
            # Store chunks in session state
            st.session_state.documents = chunks
            
            # Create FAISS index and embeddings
            index, embeddings = create_vector_store(chunks)
            st.session_state.vector_store = index
            st.session_state.embeddings = embeddings
            
            # Show success or warning
            if index is not None:
                st.success(f"✨ Created {len(chunks)} searchable chunks!")
            else:
                st.warning("No valid chunks created for embeddings.")
        else:
            st.warning("No valid materials found.")

    st.markdown("---")
    
    # Clear chat button
    if st.button("🗑️ Clear Chat"):
        st.session_state.messages = []
        st.rerun()

    # Show current status
    st.markdown("**📊 Status**")
    if st.session_state.vector_store:
        st.info(f"✅ {len(st.session_state.documents)} chunks loaded")
    else:
        st.warning("⚠️ No materials loaded")

# ============================================================
# CHAT INTERFACE
# ============================================================

st.markdown("---")

# Display chat history
for msg in st.session_state.messages:
    with st.chat_message(msg["role"]):
        st.markdown(msg["content"])

# Chat input box
if prompt := st.chat_input("Ask me anything about your materials..."):
    # Add user message to history
    st.session_state.messages.append({"role": "user", "content": prompt})
    
    # Display user message
    with st.chat_message("user"):
        st.markdown(prompt)

    # Generate and display assistant response
    with st.chat_message("assistant"):
        with st.spinner("Thinking..."):
            # Try to find relevant chunks from uploaded materials
            chunks = retrieve_relevant_chunks(prompt)
            
            if chunks and len(chunks) > 0:
                # RAG mode: Use retrieved chunks
                response = generate_response(prompt, chunks)
            else:
                # Fallback to web search
                st.info("🔍 Searching the web for information...")
                web_results = web_search(prompt)
                response = generate_response(prompt, None, web_results)
            
            # Display response
            st.markdown(response)

    # Add assistant response to history
    st.session_state.messages.append({"role": "assistant", "content": response})

# ============================================================
# FOOTER
# ============================================================

st.markdown("---")
st.markdown(
    "<div style='text-align: center; color: #666;'>"
    "Developed by <strong>Akash Bauri</strong> | "
    "Powered by <strong>Gemini (Hybrid RAG + Web)</strong>"
    "</div>",
    unsafe_allow_html=True
)
